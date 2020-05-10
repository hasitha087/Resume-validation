import PyPDF2
import os
import pickle
import shutil
import textract
from pptx import Presentation
from odf import text, teletype
from odf.opendocument import load
import zipfile
import patoolib
from datetime import datetime
import calendar
import string
import re
from nltk.stem.wordnet import WordNetLemmatizer
from nltk.stem import PorterStemmer
from bs4 import BeautifulSoup as bs
import params


d = datetime.utcnow()
unixtime = calendar.timegm(d.utctimetuple())

lemmatizer = WordNetLemmatizer()
port = PorterStemmer()


def exclude_string(text, file):
    vectorizer, classifier = model_opener()

    if any(s in text.lower() for s in params.exclude_list):
        print("Excluded by keywords: " + file)
        if not os.path.exists(params.FILTER_PATH + "noncv"):
            os.makedirs(params.FILTER_PATH + "noncv")
        shutil.copy(params.CVPATH + file, params.FILTER_PATH + "noncv")

    else:
        # text = preprocessor(text)
        sample_resume = vectorizer.transform([text]).toarray()
        prediction = classifier.predict(sample_resume)
        prob = classifier.predict_proba(sample_resume)

        # Write prediction probability into log file
        log_writer(file, prob)

        file_seperator(prob, prediction, file)


def model_opener():
    ##Open Vectorizer model
    with open(params.PKL_VECTOR, 'rb') as file:
        vectorizer = pickle.load(file)

    ##Open Classifier model
    with open(params.PKL_MODEL, 'rb') as file:
        classifier = pickle.load(file)

    return vectorizer, classifier


def preprocessor(each_resume):
    each_resume = each_resume.translate(str.maketrans('', '', string.punctuation))
    each_resume = each_resume.translate(str.maketrans('', '', "\n"))
    #each_resume = each_resume.translate(str.maketrans('', '', "\t"))
    #each_resume = each_resume.translate(str.maketrans('', '', " "))
    each_resume = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\xff]', ' ', each_resume)
    each_resume = lemmatizer.lemmatize(each_resume)
    each_resume = port.stem(each_resume)

    return each_resume


def file_seperator(prob, prediction, file):
    if prob[:, 1][0] >= 0.6:
        print("This is Resume: " + file + " " + str(round(prob[:, 1][0]*100, 2)) + "%")
        if not os.path.exists(params.FILTER_PATH + "cv"):
            os.makedirs(params.FILTER_PATH + "cv")
        shutil.copy(params.CVPATH + file, params.FILTER_PATH + "cv")
    else:
        print("This is Non-Resume: " + file + " " + str(round(prob[:, 1][0]*100, 2)) + "%")
        if not os.path.exists(params.FILTER_PATH + "noncv"):
            os.makedirs(params.FILTER_PATH + "noncv")
        shutil.copy(params.CVPATH + file, params.FILTER_PATH + "noncv")


def log_writer(file, pred_prob):
    with open(params.LOG_PATH + 'prob_{}.txt'.format(str(unixtime)), "a") as myfile:
        myfile.write(file + ";" + str(round(pred_prob[:, 1][0]*100, 2)) + "\n")


def textract_reader(filename):
    text = textract.process(filename)
    return text.decode('utf-8')


def pdf_reader(CVPATH, file):

    try:
        f = open(CVPATH + file, "rb")
        pdfReader = PyPDF2.PdfFileReader(f)
        num_pages = pdfReader.numPages
        count = 0
        text = ""
        # The while loop will read each page
        while count < num_pages:
            pageObj = pdfReader.getPage(count)
            count += 1
            text += pageObj.extractText()

        exclude_string(text, file)

    except Exception as e:
        print("Unable to handle the " + file)
        if not os.path.exists(params.FILTER_PATH + "exception"):
            os.makedirs(params.FILTER_PATH + "exception")
        shutil.copy(CVPATH + file, params.FILTER_PATH + "exception")


def docx_odt_reader(CVPATH, file):

    try:
        alltext = ' '
        word = textract_reader(CVPATH + file)
        each_resume = " "
        for j in range(len(word)):
            each_resume += word[j]

        exclude_string(each_resume, file)


    except Exception as e:
        print("Unable to handle the " + file)
        if not os.path.exists(params.FILTER_PATH + "exception"):
            os.makedirs(params.FILTER_PATH + "exception")
        shutil.copy(CVPATH + file, params.FILTER_PATH + "exception")


def doc_opener(CVPATH, file):

    try:
        soup = bs(open(CVPATH + file, encoding="Latin-1").read(), features="lxml")
        [s.extract() for s in soup(['style', 'script'])]
        word = soup.get_text()
        #word = "".join("".join(tmpText.split('\t')).split('\n')).encode('utf-8').strip()
        each_resume = " "
        for j in range(len(word)):
            each_resume += word[j]

        exclude_string(each_resume, file)


    except Exception as e:
        print("Unable to handle the " + file)
        if not os.path.exists(params.FILTER_PATH + "exception"):
            os.makedirs(params.FILTER_PATH + "exception")

        shutil.copy(CVPATH + file, params.FILTER_PATH + "exception")


def pptx_opener(CVPATH, file):

    try:
        f = open(CVPATH + file, "rb")
        prs = Presentation(f)

        each_resume = " "
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    word = shape.text

                for j in range(len(word)):
                    each_resume += word[j]

        exclude_string(each_resume, file)


    except Exception as e:
        print("Unable to handle the " + file)
        if not os.path.exists(params.FILTER_PATH + "exception"):
            os.makedirs(params.FILTER_PATH + "exception")
        shutil.copy(CVPATH + file, params.FILTER_PATH + "exception")


def odp_reader(CVPATH, file):
    vectorizer, classifier = model_opener()
    try:
        textdoc = load(CVPATH + file)
        allparas = textdoc.getElementsByType(text.P)
        each_resume = ' '
        for i in range(len(allparas)):
            word = teletype.extractText(allparas[i])
            for j in range(len(word)):
                each_resume += word[j]

        exclude_string(each_resume, file)


    except Exception as e:
        print("Unable to handle the " + file)
        if not os.path.exists(params.FILTER_PATH + "exception"):
            os.makedirs(params.FILTER_PATH + "exception")
        shutil.copy(CVPATH + file, params.FILTER_PATH + "exception")


def zipdir_reader(EXTRACT_PATH):
        # Remove all the files in zipextract dir
        for file in os.scandir(EXTRACT_PATH):
            os.unlink(file.path)

        for r, d, f in os.walk(EXTRACT_PATH):
            for file in f:
                # Read pdf files
                if str(file).split(".")[-1].lower() == "pdf":
                    pdf_reader(EXTRACT_PATH, file)

                # Read docx and odt files
                elif str(file).split(".")[-1].lower() == "docx" or str(file).split(".")[1].lower() == "odt":
                    docx_odt_reader(EXTRACT_PATH, file)

                # If file is doc move it into temp_data folder and convert into docx. Then process and remove that files from temp dir
                elif str(file).split(".")[-1].lower() == "docm" or \
                        str(file).split(".")[-1].lower() == "dot" or \
                        str(file).split(".")[-1].lower() == "dotx":
                    doc_opener(EXTRACT_PATH, file)

                # Read pptx files
                elif str(file).split(".")[-1].lower() == "pptx":
                    pptx_opener(EXTRACT_PATH, file)

                elif str(file).split(".")[-1].lower() == "odp":
                    odp_reader(EXTRACT_PATH, file)

                else:
                    if not os.path.exists(params.FILTER_PATH + "noncv"):
                        os.makedirs(params.FILTER_PATH + "noncv")
                    shutil.copy(EXTRACT_PATH + file, params.FILTER_PATH + "noncv")


def copy_sub_dir_files(EXTRACT_PATH):
    for root, dirs, files in os.walk(EXTRACT_PATH):
        if root != EXTRACT_PATH:
            for root, dirs, files in os.walk(root):
                for file in files:
                    print(root + "/" + file)
                    shutil.copy(root + "/" + file, EXTRACT_PATH)
                shutil.rmtree(root)


def main():

    try:
        ##Read Files by looking at the extension
        for r, d, f in os.walk(params.CVPATH):
            for file in f:
                # Read pdf files
                if str(file).split(".")[-1].lower() == "pdf":
                    pdf_reader(params.CVPATH, file)

                #Read docx and odt files
                elif str(file).split(".")[-1].lower() == "docx" or str(file).split(".")[1].lower() == "odt":
                    docx_odt_reader(params.CVPATH, file)

                # If file is doc move it into temp_data folder and convert into docx. Then process and remove that files from temp dir
                elif str(file).split(".")[-1].lower() == "doc" or \
                        str(file).split(".")[-1].lower() == "docm" or \
                        str(file).split(".")[-1].lower() == "dot" or \
                        str(file).split(".")[-1].lower() == "dotx":
                    doc_opener(params.CVPATH, file)

                # Read pptx files
                elif str(file).split(".")[-1].lower() == "pptx":
                    pptx_opener(params.CVPATH, file)

                elif str(file).split(".")[-1].lower() == "odp":
                    odp_reader(params.CVPATH, file)

                # Zip extractor
                elif str(file).split(".")[-1].lower() == "zip":
                    with zipfile.ZipFile(params.CVPATH + file, 'r') as zip_ref:
                        zip_ref.extractall(params.EXTRACT_PATH)
                    copy_sub_dir_files(params.EXTRACT_PATH)
                    zipdir_reader(params.EXTRACT_PATH)

                elif str(file).split(".")[-1].lower() == "rar":
                    patoolib.extract_archive(params.CVPATH + file, outdir=params.EXTRACT_PATH)
                    copy_sub_dir_files(params.EXTRACT_PATH)
                    zipdir_reader(params.EXTRACT_PATH)

                else:
                    if not os.path.exists(params.FILTER_PATH + "noncv"):
                        os.makedirs(params.FILTER_PATH + "noncv")
                    shutil.copy(params.CVPATH + file, params.FILTER_PATH + "noncv")

    except Exception as e:
        print(str(e))


if __name__ == '__main__':
    main()